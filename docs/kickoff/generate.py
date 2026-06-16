# -*- coding: utf-8 -*-
"""Generiert Word-Dokumentation und PowerPoint fuer den ALZ-Kickoff."""
import datetime
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = "/home/user/azure-landing-zone/docs/kickoff"
DATE = "16.06.2026"
VERSION = "1.1"

# Farbpalette
BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK = RGBColor(0x24, 0x3A, 0x5E)
GREY = RGBColor(0x60, 0x5E, 0x5C)
PBLUE = PRGB(0x00, 0x78, 0xD4)
PDARK = PRGB(0x24, 0x3A, 0x5E)
PLIGHT = PRGB(0xF2, 0xF6, 0xFB)
PWHITE = PRGB(0xFF, 0xFF, 0xFF)
PGREY = PRGB(0x60, 0x5E, 0x5C)

# =====================================================================
# WORD-DOKUMENTATION
# =====================================================================

def shade_cell(cell, hex_color):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear'); shd.set(qn('w:color'), 'auto'); shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)

def set_cell_text(cell, text, bold=False, color=None, size=9.5, white=False):
    cell.text = ""
    p = cell.paragraphs[0]; p.paragraph_format.space_after = Pt(2); p.paragraph_format.space_before = Pt(2)
    r = p.add_run(text); r.bold = bold; r.font.size = Pt(size)
    if white: r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    elif color: r.font.color.rgb = color

def add_table(doc, headers, rows, widths=None):
    t = doc.add_table(rows=1, cols=len(headers)); t.alignment = WD_TABLE_ALIGNMENT.CENTER
    t.style = 'Light Grid Accent 1'
    for i, h in enumerate(headers):
        set_cell_text(t.rows[0].cells[i], h, bold=True, white=True)
        shade_cell(t.rows[0].cells[i], '0078D4')
    for row in rows:
        cells = t.add_row().cells
        for i, val in enumerate(row):
            set_cell_text(cells[i], val)
    if widths:
        for i, w in enumerate(widths):
            for row in t.rows:
                row.cells[i].width = Inches(w)
    doc.add_paragraph()
    return t

def add_toc(doc):
    p = doc.add_paragraph(); run = p.add_run()
    fld = OxmlElement('w:fldChar'); fld.set(qn('w:fldCharType'), 'begin')
    instr = OxmlElement('w:instrText'); instr.set(qn('xml:space'), 'preserve'); instr.text = 'TOC \\o "1-2" \\h \\z \\u'
    sep = OxmlElement('w:fldChar'); sep.set(qn('w:fldCharType'), 'separate')
    t = OxmlElement('w:t'); t.text = "Inhaltsverzeichnis – in Word mit Strg+A, dann F9 aktualisieren."
    end = OxmlElement('w:fldChar'); end.set(qn('w:fldCharType'), 'end')
    for e in (fld, instr, sep, t, end): run._r.append(e)

def body(doc, text):
    p = doc.add_paragraph(text); p.paragraph_format.space_after = Pt(6)
    return p

def bullet(doc, text, bold_prefix=None):
    p = doc.add_paragraph(style='List Bullet')
    if bold_prefix:
        r = p.add_run(bold_prefix + ": "); r.bold = True
    p.add_run(text)
    return p

def build_docx():
    doc = Document()
    n = doc.styles['Normal']; n.font.name = 'Calibri'; n.font.size = Pt(10.5)
    for lvl, col, sz in [('Heading 1', DARK, 16), ('Heading 2', BLUE, 13), ('Heading 3', DARK, 11.5)]:
        st = doc.styles[lvl]; st.font.color.rgb = col; st.font.name = 'Calibri'; st.font.size = Pt(sz)

    # Titelseite
    for _ in range(3): doc.add_paragraph()
    t = doc.add_paragraph(); t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("Azure Landing Zone"); r.bold = True; r.font.size = Pt(30); r.font.color.rgb = BLUE
    s = doc.add_paragraph(); s.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = s.add_run("Technische Dokumentation & Architektur"); r.font.size = Pt(15); r.font.color.rgb = DARK
    s2 = doc.add_paragraph(); s2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = s2.add_run("Bicep / Azure Verified Modules · Hub-and-Spoke"); r.font.size = Pt(11); r.font.color.rgb = GREY
    for _ in range(8): doc.add_paragraph()
    m = doc.add_paragraph(); m.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = m.add_run(f"Stand: {DATE} · Version {VERSION} · Vorbereitung Kickoff"); r.font.size = Pt(10); r.font.color.rgb = GREY
    doc.add_page_break()

    doc.add_heading("Inhaltsverzeichnis", level=1)
    add_toc(doc)
    doc.add_page_break()

    # 1 Management Summary
    doc.add_heading("1. Management Summary", level=1)
    body(doc, "Dieses Dokument beschreibt die Azure Landing Zone (ALZ) – das standardisierte, "
              "richtlinien­gesteuerte Fundament für den Betrieb von Workloads in Azure. Die Umsetzung basiert "
              "auf dem offiziellen Microsoft ALZ Bicep Accelerator (Starter-Modul der öffentlichen Microsoft-"
              "Registry, ausgerollt über das ALZ-PowerShell-Modul mit Deploy-Accelerator). Damit wird das "
              "vollständige, von Microsoft gepflegte ALZ-Policy- und Governance-Set übernommen, statt eines "
              "Eigenbaus.")
    body(doc, "Die Landing Zone liefert von Beginn an eine konsistente Management-Group-Hierarchie (12 "
              "Management Groups), das volle ALZ-Policy-Set (149 Policy-Definitionen, 42 Initiativen, 123 "
              "konkrete Policy-Assignments, 5 Custom-Rollen), zentrale Protokollierung über das AMA-Pattern "
              "(Azure Monitor Agent), Netzwerk-Konnektivität nach dem Hub-and-Spoke-Muster mit Azure Firewall "
              "sowie eine Security-Baseline mit Microsoft Defender for Cloud – Letztere wird automatisch über "
              "die Policy-Assignments onboardet.")
    body(doc, "Der Aufbau erfolgt vollständig als Infrastructure-as-Code (Bicep, Azure Verified Modules) und "
              "wird über generierte CI/CD-Pipelines mit passwortloser OIDC-Anmeldung ausgerollt. Ein Bootstrap "
              "(Phase 0) richtet Repository, Identität und Pipelines automatisch ein.")
    body(doc, "Ziel des Kickoffs ist es, den technischen Ansprechpartnern des Kunden den Aufbau, den "
              "aktuellen Umsetzungsstand und die noch zu treffenden Entscheidungen vorzustellen (siehe "
              "Kapitel 10).")

    # 2 Ziele & Prinzipien
    doc.add_heading("2. Zielsetzung & Design-Prinzipien", level=1)
    for bp, tx in [
        ("Richtliniengesteuerte Governance", "Leitplanken (Guardrails) werden zentral per Azure Policy auf Management-Group-Ebene erzwungen, nicht pro Ressource."),
        ("Subscription-Demokratisierung", "Die Subscription ist die Einheit der Skalierung; Workloads werden in dedizierte Landing-Zone-Subscriptions platziert."),
        ("Einheitliche Steuerungsebene", "Eine konsistente Hierarchie und zentrales Logging/Monitoring über alle Subscriptions hinweg."),
        ("Sicherheit von Grund auf", "Defender for Cloud, Netzwerksegmentierung über Azure Firewall, private Namensauflösung und Zero-Trust-Prinzipien."),
        ("Automatisierung & Wiederholbarkeit", "Alles als Code (Bicep), versioniert, per Pipeline reproduzierbar und über What-If vorab prüfbar."),
        ("Erweiterbarkeit", "Modularer Aufbau; einzelne Domänen lassen sich unabhängig erweitern (z. B. ExpressRoute, Virtual WAN, Sentinel)."),
    ]:
        bullet(doc, tx, bp)

    # 3 Architektur
    doc.add_heading("3. Architekturüberblick", level=1)
    body(doc, "Die Landing Zone folgt der offiziellen Microsoft-ALZ-Referenz. Die Management-Group-Hierarchie "
              "unterhalb der Tenant Root strukturiert Plattform- und Workload-Subscriptions:")
    mono = doc.add_paragraph()
    rr = mono.add_run(
        "Tenant Root\n"
        "└─ alz  (Intermediate Root)\n"
        "   ├─ alz-platform\n"
        "   │   ├─ alz-platform-connectivity   → Hub-Netzwerk, Firewall, DNS\n"
        "   │   ├─ alz-platform-identity       → Identity-Dienste\n"
        "   │   ├─ alz-platform-management     → Logging, Monitoring\n"
        "   │   └─ alz-platform-security       → Security-Tooling\n"
        "   ├─ alz-landingzones\n"
        "   │   ├─ alz-landingzones-corp       → interne Workloads\n"
        "   │   ├─ alz-landingzones-online     → internetseitige Workloads\n"
        "   │   └─ alz-landingzones-local      → souveräne/vertrauliche Workloads\n"
        "   ├─ alz-sandbox                     → Experimente (lockere Policies)\n"
        "   └─ alz-decommissioned              → Stilllegung (gesperrt)")
    rr.font.name = 'Consolas'; rr.font.size = Pt(9); rr.font.color.rgb = DARK
    body(doc, "Primärregion ist Germany West Central, Sekundärregion North Europe. Plattform-Dienste laufen in "
              "den Subscriptions Management (Logging) und Connectivity (Netzwerk); Identity ist vorgesehen.")

    # 4 Domänen
    doc.add_heading("4. Domänen im Detail", level=1)

    doc.add_heading("4.1 Governance (Management Groups, Policy, RBAC)", level=2)
    body(doc, "Die Management-Group-Hierarchie und das Policy-Set stammen aus dem offiziellen Microsoft-"
              "Accelerator-Modul (avm/ptn/alz/empty, Version 0.3.6). Damit wird nicht ein schlankes Eigenbau-"
              "Set, sondern das vollständige, von Microsoft gepflegte ALZ-Policy-Set ausgerollt:")
    for bp, tx in [
        ("Policy-Definitionen", "149 Custom-Definitionen (u. a. Monitoring, Network, Storage, SQL, Guardrails je Dienst)."),
        ("Initiativen", "42 Policy-Set-Definitionen (z. B. Deploy-MDFC-Config, Deploy-Private-DNS-Zones, Enforce-Guardrails-* je Dienst, Enforce-ACSB)."),
        ("Policy-Assignments", "123 konkrete Zuweisungen, auf den passenden MG-Ebenen verankert (zzgl. Vererbung an Kind-MGs)."),
        ("Custom-Rollen", "5 Custom Role Definitions (Subscription-Owner, Security-Operations, Network-Management, Application-Owners, Network-Subnet-Contributor)."),
    ]:
        bullet(doc, tx, bp)
    body(doc, "Verteilung der direkt zugewiesenen Policy-Assignments je Management-Group-Ebene:")
    add_table(doc, ["MG-Ebene", "Assignments", "Repräsentative Beispiele"], [
        ["alz (Intermediate Root)", "17", "Deploy-MDFC-Config-H224, Deploy-MDEndpoints, Deploy-AzActivity-Log, Deploy-Diag-LogsCat, Enforce-ACSB"],
        ["landingzones", "53", "Deny-Storage-http, Deny-MgmtPorts-Internet, Deploy-VM-Monitoring, Enforce-Guardrails-* (Audit/DoNotEnforce)"],
        ["landingzones-corp", "5", "Deny-Public-Endpoints, Deny-Public-IP-On-NIC, Deny-HybridNetworking, Audit-PeDnsZones, Deploy-Private-DNS-Zones"],
        ["platform", "40", "Deploy-VM-Monitoring, Enforce-Backup, Enforce-Guardrails-* (Audit), DenyAction-DeleteUAMIAMA"],
        ["platform-identity", "4", "Deny-MgmtPorts-Internet, Deny-Public-IP, Deny-Subnet-Without-Nsg, Deploy-VM-Backup"],
        ["platform-connectivity", "1", "Enable-DDoS-VNET"],
        ["landingzones-local", "1", "Enforce-ALDO-Services (Azure Local disconnected)"],
        ["sandbox", "1", "Enforce-ALZ-Sandbox (gelockerte Guardrails)"],
        ["decommissioned", "1", "Enforce-ALZ-Decomm (Stilllegung)"],
    ], widths=[2.0, 1.0, 4.0])
    body(doc, "Viele Assignments sind vom Typ DeployIfNotExists oder Modify: sie legen bei Bedarf automatisch "
              "Remediation-Identities (Managed Identities) an und konfigurieren Ziel-Ressourcen nach (z. B. "
              "Diagnose-Einstellungen, Defender-Onboarding, Monitoring-Agents). Die Guardrails der Initiativen "
              "Enforce-Guardrails-* sind im Auslieferungszustand teils auf DoNotEnforce gesetzt (auditierend) und "
              "können pro Compliance-Bedarf scharfgeschaltet werden.")
    body(doc, "RBAC: Das Accelerator-Modul liefert 5 Custom-Rollen und weist Built-in-/Custom-Rollen an Entra-"
              "ID-Gruppen je Management Group zu. Object-IDs werden über die Konfiguration injiziert (keine "
              "Hardcodes); ohne gesetzte Gruppen-IDs bleibt die Zuweisung ein gefahrloser No-Op.")

    doc.add_heading("4.2 Identity", level=2)
    body(doc, "Die Management Group alz-platform-identity sowie eine Identity-Subscription sind vorgesehen. "
              "Konkrete Identity-Ressourcen (z. B. Entra-ID-Diagnose an Log Analytics, hybride Anbindung, "
              "Domänendienste) sind Teil der Roadmap und werden gemäß Kundenanforderung ausgestaltet "
              "(siehe Beratungspunkte, Kapitel 10).")

    doc.add_heading("4.3 Management & Monitoring (Logging)", level=2)
    body(doc, "Die zentrale Protokollierung in der Management-Subscription wird über das offizielle AMA-Pattern "
              "(avm/ptn/alz/ama, Version 0.2.0) bereitgestellt – der Azure Monitor Agent samt zugehöriger "
              "Data Collection Rules. Eine eigene DCR-Logik entfällt:")
    for bp, tx in [
        ("Log Analytics Workspace", "law-alz-<Region>, Aufbewahrung 365 Tage, SKU PerGB2018."),
        ("Data Collection Rules", "drei DCRs: VM Insights, Change Tracking und Defender-for-SQL."),
        ("User-Assigned Managed Identity", "mi-alz-<Region> für die Datensammlung der Agents."),
        ("Solutions", "ChangeTracking im Workspace; Microsoft Sentinel per Schalter aktivierbar."),
        ("Automation Account", "Optional zuschaltbar (Update-/Konfigurationsmanagement; default aus)."),
    ]:
        bullet(doc, tx, bp)
    body(doc, "Die VM-/Monitoring-Policies des Policy-Sets (Deploy-VM-Monitoring, Deploy-VM-ChangeTrack u. a.) "
              "verbinden neue Maschinen automatisch mit dem Workspace und den DCRs.")

    doc.add_heading("4.4 Security", level=2)
    body(doc, "Microsoft Defender for Cloud wird nicht mehr über ein eigenes Template aktiviert, sondern "
              "automatisch über die Policy-Assignments auf der alz-Root-Ebene onboardet (DeployIfNotExists). Die "
              "Sicherheits-Konfiguration ist damit Teil des Governance-Sets:")
    for bp, tx in [
        ("Deploy-MDFC-Config-H224", "Onboarding der Defender-for-Cloud-Konfiguration inkl. Defender-Pläne und Security Contacts."),
        ("Deploy-MDEndpoints / Deploy-MDEndpointsAMA", "Microsoft Defender for Endpoint und dessen Integration mit Defender for Cloud."),
        ("Deploy-MDFC-OssDb / Deploy-MDFC-SqlAtp", "Advanced Threat Protection für Open-Source-Datenbanken sowie SQL-Server/-Managed-Instances."),
        ("Deploy-AzActivity-Log / Deploy-Diag-LogsCat", "Activity-Log und Diagnose-Kategorien je Subscription/Ressource zentral ins Log Analytics."),
        ("Enforce-ACSB / Deploy-ASC-Monitoring", "Azure Compute Security Baseline und Microsoft Cloud Security Benchmark."),
    ]:
        bullet(doc, tx, bp)
    body(doc, "Kostenhinweis: Die Defender-Pläne werden über die Policy-Konfiguration auf den Standard-Tier "
              "gehoben (kostenpflichtig, in der Regel pro Ressource bzw. transaktionsbasiert). Der Tier ist über "
              "die Konfiguration steuerbar. Microsoft Sentinel ist über einen Schalter am Log Analytics Workspace "
              "vorbereitet; Data Connectors und Analyseregeln nach Bedarf.")

    doc.add_heading("4.5 Connectivity (Hub-and-Spoke)", level=2)
    body(doc, "Das Hub-Netzwerk in der Connectivity-Subscription bildet das Rückgrat. Spokes (Workloads) "
              "werden an den Hub gepeert und leiten ausgehenden Verkehr über die Azure Firewall. Wichtig: Der "
              "Microsoft-Standard des Accelerators aktiviert ALLE Dienste in BEIDEN Regionen.")
    for bp, tx in [
        ("Hub-VNets", "Primär 10.0.0.0/22 (Germany West Central), Sekundär 10.1.0.0/22 (North Europe), bidirektional gepeert."),
        ("Azure Firewall + Policy", "2x Azure Firewall (Standard), je Hub eine – mit Firewall Policy."),
        ("Azure Bastion", "2x Bastion (je Hub) – sicherer RDP/SSH-Zugang ohne öffentliche IPs."),
        ("VPN Gateway", "2x VPN Gateway (VpnGw1AZ, active-active BGP, ASN 65515)."),
        ("ExpressRoute Gateway", "2x ExpressRoute Gateway (je Hub)."),
        ("DDoS Network Protection", "Aktiviert (Hub 1) – teuerster Einzelposten (~€2.500/Monat)."),
        ("DNS Private Resolver & Private DNS Zones", "2x DNS Private Resolver, zentrale Private DNS Zones für Private Endpoints."),
    ]:
        bullet(doc, tx, bp)
    body(doc, "Kostenhinweis: Im Microsoft-Default summieren sich diese Dienste auf rund €5.800/Monat (doppelte "
              "Firewalls, Bastion, VPN-/ExpressRoute-Gateways plus DDoS). Für einen kostenarmen Erst-Rollout gibt "
              "es zwei Wege: network_type: none (nur Management Groups + Policies + Logging, ≈ €0) oder die "
              "deploy*-Schalter je Dienst in templates/networking/hubnetworking/main.bicepparam auf false setzen "
              "(z. B. nur VNets + Private DNS Zones, ≈ €15/Monat). Details im Bootstrap-Runbook "
              "(docs/ACCELERATOR-BOOTSTRAP.md).")

    doc.add_heading("4.6 Landing Zones & Subscription-Platzierung", level=2)
    body(doc, "Workload-Landing-Zones erhalten ein Spoke-VNet mit Route Table (Default-Route über die Firewall), "
              "bidirektionalem Hub-Peering und Verknüpfung zu den zentralen Private-DNS-Zonen. Die Platzierung "
              "von Subscriptions erfolgt über die vom Accelerator erstellte Management-Group-Struktur und den "
              "Bootstrap: bestehende Subscriptions werden in die passende Management Group (corp/online/local "
              "bzw. Plattform) eingehängt und erben sofort deren Policies und RBAC. Ein dediziertes Subscription-"
              "Vending-Template wird nicht mehr aktiv betrieben (die Vorgänger-Eigenimplementierung liegt "
              "archiviert unter legacy-custom/).")

    # 5 Ressourcen-Übersicht
    doc.add_heading("5. Erstellte Azure-Ressourcen – Gesamtübersicht", level=1)
    body(doc, "Die folgende Übersicht fasst die Ressourcen zusammen, die der ALZ Bicep Accelerator erstellt. "
              "Governance-Objekte (Management Groups, Policies, Rollen) liegen auf Tenant-/MG-Scope, alle "
              "übrigen werden in den Platform-Subscriptions angelegt. Insgesamt nutzt das Starter-Modul 22 "
              "Azure Verified Modules (u. a. avm/ptn/alz/empty, avm/ptn/alz/ama, Netzwerk-Module).")

    doc.add_heading("5.1 Governance (Tenant-/MG-Scope)", level=2)
    add_table(doc,
        ["Ressourcentyp", "Name / Konvention", "Anzahl", "Kosten"],
        [
            ["Management Group", "alz, alz-platform, alz-platform-{connectivity,identity,management,security}, "
             "alz-landingzones, alz-landingzones-{corp,online,local}, alz-sandbox, alz-decommissioned", "12", "Kostenlos"],
            ["Policy Definition (Custom)", "vollständiges ALZ-Set, geladen aus lib/alz/", "149", "Kostenlos"],
            ["Policy Set Definition (Initiative)", "Deploy-MDFC-Config, Deploy-Private-DNS-Zones, Enforce-Guardrails-* …", "42", "Kostenlos"],
            ["Policy Assignment", "auf alz / landingzones / corp / platform / identity … verteilt", "123", "Kostenlos"],
            ["Custom Role Definition", "Subscription-Owner, Security-Operations, Network-Management, Application-Owners, Network-Subnet-Contributor", "5", "Kostenlos"],
            ["Remediation-Identity (auto)", "Managed Identities für DeployIfNotExists/Modify", "n", "Kostenlos"],
            ["RBAC Role Assignment", "Built-in-/Custom-Rollen je Entra-Gruppe (konfigurierbar)", "n", "Kostenlos"],
        ],
        widths=[2.3, 3.6, 0.7, 1.4])

    doc.add_heading("5.2 Logging (Management-Subscription, avm/ptn/alz/ama)", level=2)
    add_table(doc,
        ["Ressourcentyp", "Name / Konvention", "Anzahl", "Kosten"],
        [
            ["Resource Group", "rg-alz-logging-<region>", "1", "Kostenlos"],
            ["Log Analytics Workspace", "law-alz-<region> (365 Tage, PerGB2018)", "1", "5 GB/Mon. kostenlos"],
            ["Data Collection Rule – VM Insights", "dcr-vmi-alz-<region>", "1", "Im LAW enthalten"],
            ["Data Collection Rule – Change Tracking", "dcr-ct-alz-<region>", "1", "Im LAW enthalten"],
            ["Data Collection Rule – Defender SQL", "dcr-mdfcsql-alz-<region>", "1", "Im LAW enthalten"],
            ["User-Assigned Managed Identity", "mi-alz-<region>", "1", "Kostenlos"],
            ["LAW Solution – ChangeTracking", "ChangeTracking", "1", "Im LAW enthalten"],
            ["Automation Account (optional, default aus)", "aa-alz-<region> (Basic)", "0–1", "Verbrauchsabhängig"],
        ],
        widths=[2.3, 3.6, 0.7, 1.4])

    doc.add_heading("5.3 Security – Defender for Cloud (via Policy)", level=2)
    body(doc, "Microsoft Defender for Cloud wird über die Deploy-MDFC-*-Policy-Assignments auf der alz-Root-"
              "Ebene automatisch onboardet (DeployIfNotExists) – kein eigenes Security-Template mehr. Die "
              "Defender-Pläne werden auf den Standard-Tier gehoben (kostenpflichtig); der Tier ist konfigurierbar.")
    add_table(doc,
        ["Plan / Komponente", "aktiviert über", "Kosten (Standard)"],
        [
            ["Defender-Konfiguration + Security Contacts", "Deploy-MDFC-Config-H224", "siehe Pläne"],
            ["Defender for Servers (VMs)", "Deploy-MDFC-Config-H224", "~€13/Server/Mon."],
            ["Defender for Endpoint", "Deploy-MDEndpoints / Deploy-MDEndpointsAMA", "im Server-Plan"],
            ["Defender for SQL (Server/MI/VM)", "Deploy-MDFC-SqlAtp / Deploy-MDFC-DefSQL-AMA", "~€13/Server/Mon."],
            ["Defender for Open-Source-DB", "Deploy-MDFC-OssDb", "~€13/Server/Mon."],
            ["Microsoft Cloud Security Benchmark / ACSB", "Deploy-ASC-Monitoring / Enforce-ACSB", "Kostenlos (Audit)"],
            ["Activity-Log / Diagnose → LAW", "Deploy-AzActivity-Log / Deploy-Diag-LogsCat", "Im LAW enthalten"],
        ],
        widths=[2.6, 3.0, 1.4])

    doc.add_heading("5.4 Hub Networking (Connectivity-Subscription)", level=2)
    body(doc, "Achtung: Im Microsoft-Default sind alle Dienste in BEIDEN Hubs aktiviert (Summe ≈ €5.800/Monat). "
              "Die Anzahl-Spalte zeigt den Default; über die deploy*-Schalter je Dienst lässt sich reduzieren.")
    add_table(doc,
        ["Ressourcentyp", "Name / Konvention", "Anzahl", "Kosten (Default)"],
        [
            ["Resource Group", "rg-alz-conn-<region>", "2", "Kostenlos"],
            ["Virtual Network (Hub 1 / Hub 2)", "10.0.0.0/22 (GWC) · 10.1.0.0/22 (NE)", "2", "Kostenlos"],
            ["Azure Firewall (Standard)", "je Hub eine", "2", "~€1.100/Hub/Monat"],
            ["Firewall Policy", "je Hub eine", "2", "Kostenlos"],
            ["Azure Bastion", "je Hub eine", "2", "~€120/Hub/Monat"],
            ["VPN Gateway", "VpnGw1AZ, active-active BGP, ASN 65515", "2", "~€140/Hub/Monat"],
            ["ExpressRoute Gateway", "je Hub eines", "2", "~€280/Hub/Monat"],
            ["DDoS Network Protection", "Hub 1", "1", "~€2.500/Monat"],
            ["DNS Private Resolver", "Inbound/Outbound-Subnetze", "2", "~€25/Hub/Monat"],
            ["Private DNS Zones", "privatelink.* für Private Endpoints", "viele", "~€0,90/Zone/Mon."],
            ["VNet Peering Hub↔Hub", "Hub 1 ↔ Hub 2 (bidirektional)", "2", "Gering (Datenübertragung)"],
        ],
        widths=[2.3, 3.2, 0.7, 1.8])
    body(doc, "Kostenarme Variante: network_type: none deployt nur Management Groups + volles Policy-Set + "
              "Logging (≈ €0); alternativ je Dienst die deploy*-Schalter in der main.bicepparam auf false setzen.")

    doc.add_heading("5.5 Spoke Networking (Workload-Subscription, je Landing Zone)", level=2)
    add_table(doc,
        ["Ressourcentyp", "Name / Konvention", "Anzahl", "Kosten"],
        [
            ["Resource Group", "rg-alz-spoke-<workload>-<region>", "1", "Kostenlos"],
            ["Virtual Network (Spoke)", "vnet-<workload>-<region> (z. B. 10.2.0.0/24)", "1", "Kostenlos"],
            ["Route Table", "rt-<workload>-<region> (Default-Route → Firewall)", "1", "Kostenlos"],
            ["VNet Peering Spoke→Hub", "spoke-to-hub", "1", "Gering (Datenübertragung)"],
            ["VNet Peering Hub→Spoke", "hub-to-spoke", "1", "Gering (Datenübertragung)"],
            ["Private DNS Zone Link", "Link je Zone zu Spoke VNet", "viele", "Kostenlos"],
        ],
        widths=[2.3, 3.2, 0.7, 1.8])

    doc.add_heading("5.6 Subscription-Platzierung (MG-Scope)", level=2)
    body(doc, "Der Accelerator platziert Subscriptions über die Management-Group-Struktur und den Bootstrap; ein "
              "eigenes Subscription-Vending-Template wird nicht mehr aktiv betrieben (archiviert unter "
              "legacy-custom/).")
    add_table(doc,
        ["Modus", "Aktion", "Voraussetzung", "Kosten"],
        [
            ["Placement (Standard)", "Bestehende Subscription in Ziel-MG einhängen → erbt Policies/RBAC", "Keine Billing-Rechte nötig", "Kostenlos"],
            ["Neuanlage (optional)", "Neue Subscription per EA/MCA, danach in die Ziel-MG", "Billing Account Owner-Rolle", "Sub-abhängig"],
        ],
        widths=[1.8, 3.3, 2.4, 0.5])

    # 6 IaC & CI/CD
    doc.add_heading("6. Infrastructure as Code & CI/CD", level=1)
    body(doc, "Der Rollout erfolgt über das offizielle ALZ-PowerShell-Modul (Deploy-Accelerator), nicht mehr über "
              "ein eigenes deploy.ps1 oder einen handgeschriebenen Workflow:")
    for bp, tx in [
        ("Bicep + AVM", "Deklarative Templates des Starter-Moduls auf Basis von 22 Azure Verified Modules aus der Microsoft-Registry."),
        ("Bootstrap (Phase 0)", "Terraform-basiert; erstellt das GitHub-Repository, passwortloses OIDC (Managed Identity + Federated Credentials) und die Deployment-Pipelines."),
        ("Deploy-Accelerator", "Liest config/inputs-github.yaml (Bootstrap) und config/platform-landing-zone.yaml (Starter-Modul) und richtet alles ein."),
        ("18 geordnete Deploy-Stufen", "Gemäß .config/ALZ-Powershell.config.json: Governance/MGs (1–15, inkl. Cross-MG-RBAC), Core-Logging (16), Networking Hub bzw. Virtual WAN (17/18)."),
        ("Generierte Pipelines", "Die im neuen Repo erzeugten Pipelines fahren die Stufen aus; What-If vorab, Apply über ein Approval-Gate."),
        ("Passwortlose Anmeldung", "OIDC Federated Identity ohne gespeicherte Secrets/Zertifikate."),
    ]:
        bullet(doc, tx, bp)
    body(doc, "Das vollständige Runbook (Voraussetzungen, Single-Subscription-Setup, kostenarme Variante, "
              "Aufräumen) steht in docs/ACCELERATOR-BOOTSTRAP.md.")

    # 7 Naming & IP
    doc.add_heading("7. Namens- & Adresskonzept", level=1)
    add_table(doc, ["Element", "Konvention / Bereich"], [
        ["Management Groups", "alz, alz-platform-*, alz-landingzones-*"],
        ["Resource Groups", "rg-alz-<zweck>-<region>"],
        ["Ressourcen", "law-/mi-/afw-/bas-/vnet-/afwp-alz-<region>"],
        ["Hub primär / sekundär", "10.0.0.0/22 / 10.1.0.0/22"],
        ["Spoke (Beispiel)", "10.2.0.0/24 (überlappungsfrei zu Hubs/On-Prem)"],
        ["Regionen", "germanywestcentral (primär), northeurope (sekundär)"],
    ], widths=[2.4, 4.6])

    # 8 Betrieb
    doc.add_heading("8. Betrieb & Verifikation", level=1)
    body(doc, "Der Rollout führt – nach Risiko/Kosten geordnet – von der lokalen Validierung über den Bootstrap "
              "und What-If bis zum freigegebenen Apply über die generierten Pipelines:")
    add_table(doc, ["Stufe", "Inhalt", "Azure-Wirkung"], [
        ["0 Statisch", "bicep build der Templates (lokal/CI)", "keine"],
        ["1 Bootstrap", "Deploy-Accelerator: Repo + OIDC + Pipelines (Phase 0)", "gering (State-Storage)"],
        ["2 What-If", "Vorschau je Deploy-Stufe über die Pipeline", "keine"],
        ["3 Governance/MGs", "Stufen 1–15: Hierarchie + volles Policy-Set + RBAC", "kostenlos"],
        ["4 Logging", "Stufe 16: Log Analytics + AMA/DCRs", "gering"],
        ["5 Networking", "Stufen 17/18: Hub (FW, Bastion, GW, DDoS) bzw. vWAN", "kostenrelevant"],
    ], widths=[1.7, 3.5, 1.8])
    body(doc, "Apply erfolgt jeweils erst nach Freigabe über das Approval-Gate. Für einen risikofreien Erst-"
              "Rollout kann das Netzwerk über network_type: none oder die deploy*-Schalter ausgespart werden "
              "(siehe docs/ACCELERATOR-BOOTSTRAP.md).")

    # 9 Status
    doc.add_heading("9. Umsetzungsstatus", level=1)
    body(doc, "Die Plattform basiert auf dem offiziellen Microsoft-Accelerator (Spiegel von Azure/alz-bicep-"
              "accelerator); das volle ALZ-Policy-Set ist umgesetzt:")
    add_table(doc, ["Domäne", "Status"], [
        ["Management Groups + Hierarchie (12 MGs)", "Umgesetzt"],
        ["Volles ALZ-Policy-Set (149 Defs / 42 Initiativen / 123 Assignments)", "Umgesetzt"],
        ["RBAC (5 Custom-Rollen + Zuweisungen)", "Umgesetzt (Gruppen-IDs durch Kunde)"],
        ["Logging via avm/ptn/alz/ama (LAW, 3 DCRs, MI, ChangeTracking)", "Umgesetzt"],
        ["Defender for Cloud (via Deploy-MDFC-* Policies)", "Umgesetzt (automatisches Onboarding)"],
        ["Hub-Networking, Firewall+Policy, Bastion, DNS, Peering", "Umgesetzt (Default: in beiden Regionen an)"],
        ["DNS Private Resolver", "Umgesetzt (Default an)"],
        ["VPN/ExpressRoute Gateway, DDoS", "Umgesetzt (Default an, kostenrelevant)"],
        ["Bootstrap + generierte CI/CD-Pipelines + OIDC", "Umgesetzt"],
        ["Identity-Ressourcen (über Policies hinaus)", "Roadmap"],
        ["Sentinel-Onboarding (Connectors/Regeln)", "Roadmap"],
        ["Virtual WAN (Alternative zu Hub-and-Spoke)", "Vorbereitet / optional"],
    ], widths=[4.6, 2.4])

    # 10 Roadmap
    doc.add_heading("10. Roadmap / nächste Schritte", level=1)
    for tx in [
        "Bootstrap im Kunden-Tenant ausführen (Deploy-Accelerator), dann What-If → Apply über die Pipelines.",
        "Netzwerk-Scope festlegen: voller Microsoft-Default (~€5.800/Monat) vs. kostenarme Variante (network_type=none oder deploy*-Schalter aus).",
        "Guardrail-Initiativen mit DoNotEnforce nach Compliance-Bedarf scharfschalten.",
        "Identity-Domäne ausgestalten (Entra-Diagnose, hybride Anbindung nach Bedarf).",
        "Sentinel-Onboarding mit Data Connectors und Analyseregeln.",
        "On-Prem-Konnektivität (VPN oder ExpressRoute) konfigurieren/abstimmen.",
        "Optionale Härtung: Resource Locks, erweiterte Diagnostics, Cost Management.",
    ]:
        bullet(doc, tx)

    # 11 Entscheidungspunkte
    doc.add_heading("11. Entscheidungs- & Beratungspunkte (Kickoff)", level=1)
    body(doc, "Folgende Punkte sind gemeinsam mit dem Kunden zu klären – sie bestimmen Ausprägung, Kosten und "
              "Compliance der Landing Zone:")
    add_table(doc, ["Thema", "Zu klären", "Empfehlung"], [
        ["Netzwerk-Scope & Kosten", "Voller Default (~€5.800/Mon.) vs. kostenarm", "Smoke Run: network_type=none / Schalter aus"],
        ["Netzwerk-Topologie", "Hub-and-Spoke vs. Virtual WAN", "Hub-and-Spoke (umgesetzt)"],
        ["IP-Adresskonzept", "Bereiche, Überlappung mit On-Prem", "Frühzeitig abstimmen"],
        ["Regionen / Datenresidenz", "Primär/Sekundär, Compliance", "GWC + NE bestätigen"],
        ["On-Prem-Anbindung", "VPN oder ExpressRoute (im Default beide an)", "Nach Bandbreite/SLA, sonst aus"],
        ["DDoS Protection", "Default an (~€2.500/Mon.)", "Bewusst entscheiden / für Internet-facing"],
        ["Identity", "Entra-only / hybrid / AD DS", "Anforderungsabhängig"],
        ["RBAC-Modell", "Entra-Gruppen-IDs, PIM", "Gruppen je MG definieren"],
        ["Guardrail-Enforcement", "DoNotEnforce-Initiativen scharfschalten", "Nach Compliance-Pflicht"],
        ["Defender-Tier", "Standard vs. Free je Plan", "Standard für Prod"],
        ["Sentinel", "Onboarding & Connectors", "Bei SIEM-Bedarf"],
        ["Subscriptions", "Single- vs. Multi-Subscription", "Single zum Start möglich"],
        ["Log-Aufbewahrung", "365 Tage / Kosten", "Nach Compliance"],
    ], widths=[1.8, 2.9, 2.3])

    # 12 Voraussetzungen
    doc.add_heading("12. Voraussetzungen", level=1)
    for tx in [
        "PowerShell ≥ 7.4, Azure CLI (angemeldet via az login) und das ALZ-PowerShell-Modul (Install-Module ALZ).",
        "Owner auf den vier Platform-Subscriptions (management/identity/connectivity/security – Single-Subscription möglich, alle vier IDs identisch) und erhöhter Zugriff am Tenant Root (MG-Erstellung).",
        "GitHub Personal Access Token (Scopes repo, workflow, admin:org) für den Bootstrap.",
        "Rollout über Deploy-Accelerator; Bootstrap erstellt Repo, passwortloses OIDC und Pipelines automatisch.",
        "Befüllte Konfiguration: config/inputs-github.yaml und config/platform-landing-zone.yaml (Runbook docs/ACCELERATOR-BOOTSTRAP.md).",
        "Für Subscription-Neuanlage zusätzlich: EA/MCA-Billing-Rolle.",
    ]:
        bullet(doc, tx)

    # 13 Glossar
    doc.add_heading("13. Glossar", level=1)
    add_table(doc, ["Begriff", "Bedeutung"], [
        ["ALZ", "Azure Landing Zone – standardisiertes Cloud-Fundament"],
        ["AVM", "Azure Verified Modules – geprüfte Bicep-Module von Microsoft"],
        ["MG", "Management Group – hierarchische Gruppierung von Subscriptions"],
        ["Hub-and-Spoke", "Zentraler Hub (Shared Services) + angebundene Spokes (Workloads)"],
        ["OIDC", "OpenID Connect – passwortlose Pipeline-Anmeldung an Azure"],
        ["DCR", "Data Collection Rule – Regel zur Telemetrie-Erfassung"],
        ["Vending", "Automatisierte Bereitstellung/Platzierung von Subscriptions"],
        ["What-If", "Vorschau der Änderungen ohne Deployment"],
    ], widths=[1.8, 5.2])

    # Footer
    section = doc.sections[0]
    footer = section.footer.paragraphs[0]; footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = footer.add_run("Azure Landing Zone – Technische Dokumentation · vertraulich"); fr.font.size = Pt(8); fr.font.color.rgb = GREY

    path = f"{OUT}/Azure-Landing-Zone-Dokumentation.docx"
    doc.save(path)
    return path

# =====================================================================
# POWERPOINT
# =====================================================================

def build_pptx():
    prs = Presentation()
    prs.slide_width = PInches(13.333); prs.slide_height = PInches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def add_rect(slide, x, y, w, h, color, line=False):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        if not line: shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def txt(slide, x, y, w, h, text, size, color=PDARK, bold=False, align=PP_ALIGN.LEFT, font='Calibri'):
        tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text; f = r.font
        f.size = size; f.bold = bold; f.color.rgb = color; f.name = font
        return tb

    def bullets(slide, x, y, w, h, items, size=PPt(16), color=PDARK):
        tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
        first = True
        for it in items:
            lvl = 0; t = it
            if isinstance(it, tuple): t, lvl = it
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.level = lvl
            r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + t
            r.font.size = size if lvl == 0 else PPt(13)
            r.font.color.rgb = color; r.font.name = 'Calibri'
            p.space_after = PPt(6)
        return tb

    def content_slide(title, subtitle=None):
        s = prs.slides.add_slide(blank)
        add_rect(s, 0, 0, SW, PInches(1.05), PBLUE)
        add_rect(s, 0, PInches(1.05), SW, Emu(45720), PDARK)  # thin accent line
        txt(s, PInches(0.5), PInches(0.18), PInches(12.3), PInches(0.7), title, PPt(26), PWHITE, True)
        if subtitle:
            txt(s, PInches(0.5), PInches(1.2), PInches(12.3), PInches(0.5), subtitle, PPt(14), PGREY, True)
        # footer
        txt(s, PInches(0.5), PInches(7.05), PInches(9), PInches(0.35), "Azure Landing Zone · Kickoff", PPt(9), PGREY)
        return s

    # --- Titel ---
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, PDARK)
    add_rect(s, 0, PInches(2.7), SW, PInches(0.08), PBLUE)
    txt(s, PInches(0.8), PInches(2.85), PInches(11.7), PInches(1.2), "Azure Landing Zone", PPt(48), PWHITE, True)
    txt(s, PInches(0.85), PInches(4.0), PInches(11.7), PInches(0.7), "Architektur, Umsetzungsstand & Beratung", PPt(22), PRGB(0x9F, 0xD3, 0xFF))
    txt(s, PInches(0.85), PInches(6.4), PInches(11.7), PInches(0.5), f"Kickoff · {DATE}", PPt(14), PRGB(0xB8,0xC4,0xD9))

    # --- Agenda ---
    s = content_slide("Agenda")
    bullets(s, PInches(0.7), PInches(1.5), PInches(12), PInches(5.4), [
        "Was ist eine Azure Landing Zone & Ziele",
        "Zielarchitektur und Management-Group-Hierarchie",
        "Domänen: Governance, Identity, Management, Security, Connectivity, Landing Zones",
        "Infrastructure as Code & CI/CD",
        "Umsetzungsstand & Roadmap",
        "Entscheidungs- und Beratungspunkte (gemeinsam)",
        "Nächste Schritte / Smoke Run",
    ], size=PPt(18))

    # --- Was ist ALZ ---
    s = content_slide("Was ist eine Azure Landing Zone?")
    bullets(s, PInches(0.7), PInches(1.5), PInches(12), PInches(5.4), [
        "Standardisiertes, richtliniengesteuertes Fundament für Workloads in Azure",
        "Liefert von Tag 1: Governance, Netzwerk, Identität, Security, Monitoring",
        ("Skalierungseinheit ist die Subscription – Workloads in dedizierten Landing Zones", 1),
        ("Leitplanken zentral per Azure Policy, nicht pro Ressource", 1),
        "Vollständig als Code (Bicep / Azure Verified Modules), per Pipeline reproduzierbar",
        "Ergebnis: schneller, sicherer, konsistenter Onboarding-Prozess für neue Workloads",
    ], size=PPt(17))

    # --- Prinzipien ---
    s = content_slide("Design-Prinzipien")
    bullets(s, PInches(0.7), PInches(1.5), PInches(12), PInches(5.4), [
        "Richtliniengesteuerte Governance (Guardrails auf MG-Ebene)",
        "Einheitliche Steuerungsebene & zentrales Logging",
        "Sicherheit von Grund auf (Defender, Segmentierung, Private DNS)",
        "Automatisierung & Wiederholbarkeit (IaC, What-If, CI/CD)",
        "Modulare Erweiterbarkeit je Domäne",
    ], size=PPt(18))

    # --- Architektur (Hierarchie) ---
    s = content_slide("Zielarchitektur – Management-Group-Hierarchie")
    tb = s.shapes.add_textbox(PInches(0.7), PInches(1.4), PInches(12), PInches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    tree = ("Tenant Root\n"
            "└─ alz (Intermediate Root)\n"
            "    ├─ alz-platform → connectivity · identity · management · security\n"
            "    ├─ alz-landingzones → corp · online · local\n"
            "    ├─ alz-sandbox (lockere Policies)\n"
            "    └─ alz-decommissioned (gesperrt)")
    p = tf.paragraphs[0]; r = p.add_run(); r.text = tree
    r.font.name = 'Consolas'; r.font.size = PPt(18); r.font.color.rgb = PDARK
    txt(s, PInches(0.7), PInches(6.2), PInches(12), PInches(0.6),
        "Primär: Germany West Central · Sekundär: North Europe", PPt(13), PGREY, True)

    # --- Domänen-Folien ---
    s = content_slide("Governance", "Volles ALZ-Policy-Set (offizieller Accelerator)")
    bullets(s, PInches(0.7), PInches(1.7), PInches(12), PInches(5.0), [
        "12 Management Groups + volles ALZ-Set: 149 Policy-Definitionen · 42 Initiativen · 123 Assignments · 5 Custom-Rollen",
        "alz (Root): 17 Assignments – v. a. Defender-Onboarding, Activity-Log/Diagnose, Compute Security Baseline",
        "landingzones: 53 · platform: 40 · corp: 5 (Deny Public Endpoints/IP, Hybrid Networking, Private DNS)",
        "DeployIfNotExists/Modify legen Remediation-Identities automatisch an",
        "RBAC: 5 Custom-Rollen + Zuweisung an Entra-Gruppen je MG (Object-IDs durch Kunde)",
    ], size=PPt(16))

    s = content_slide("Identity", "Plattform-Domäne (Roadmap)")
    bullets(s, PInches(0.7), PInches(1.7), PInches(12), PInches(5.0), [
        "Management Group alz-platform-identity & Identity-Subscription vorgesehen",
        "Geplant: Entra-ID-Diagnose → Log Analytics, rollenbasierte Zugriffe",
        "Optional: hybride Anbindung / Domänendienste je nach Anforderung",
        ("Ausprägung wird im Kickoff abgestimmt", 1),
    ], size=PPt(17))

    s = content_slide("Management & Monitoring", "Zentrale Protokollierung (avm/ptn/alz/ama)")
    bullets(s, PInches(0.7), PInches(1.7), PInches(12), PInches(5.0), [
        "AMA-Pattern (Azure Monitor Agent) statt eigener DCR-Logik",
        "Log Analytics Workspace (365 Tage, PerGB2018)",
        "3 Data Collection Rules: VM Insights · Change Tracking · Defender-for-SQL",
        "User-Assigned Managed Identity · ChangeTracking-Solution",
        "Automation Account & Microsoft Sentinel per Schalter zuschaltbar",
    ], size=PPt(17))

    s = content_slide("Security", "Defender for Cloud – automatisch via Policy")
    bullets(s, PInches(0.7), PInches(1.7), PInches(12), PInches(5.0), [
        "Kein eigenes Security-Template mehr – Onboarding über Deploy-MDFC-* Assignments auf alz",
        "Deploy-MDFC-Config · Deploy-MDEndpoints · Deploy-MDFC-OssDb · Deploy-MDFC-SqlAtp",
        "Activity-Log/Diagnose je Subscription → Log Analytics · Compute Security Baseline (Enforce-ACSB)",
        "Defender-Pläne auf Standard-Tier (kostenpflichtig, pro Ressource) · Sentinel vorbereitet",
    ], size=PPt(17))

    s = content_slide("Connectivity", "Hub-and-Spoke · Microsoft-Default = alles an")
    bullets(s, PInches(0.7), PInches(1.7), PInches(12), PInches(5.0), [
        "2 Hub-VNets: 10.0.0.0/22 (GWC) & 10.1.0.0/22 (NE), bidirektional gepeert",
        "Default aktiviert ALLE Dienste in BEIDEN Regionen: 2x Firewall · 2x Bastion · 2x VPN GW · 2x ER GW",
        "DDoS Network Protection (Hub 1) · 2x DNS Private Resolver · Private DNS Zones",
        "Summe Default ≈ €5.800/Monat (DDoS ~€2.500 davon)",
        "Kostenarm: network_type=none (nur MGs+Policies+Logging) oder deploy*-Schalter aus",
    ], size=PPt(16))

    s = content_slide("Landing Zones & Subscription-Platzierung")
    bullets(s, PInches(0.7), PInches(1.7), PInches(12), PInches(5.0), [
        "Spoke-VNet je Workload: Route Table → Firewall, Hub-Peering, Private-DNS-Links",
        "Segmente: corp (intern) · online (internetseitig) · local (souverän)",
        "Platzierung über die Accelerator-MG-Struktur & Bootstrap (kein eigenes Vending-Template)",
        "Bestehende Subscriptions einhängen → erben Policies/RBAC · Single-Subscription möglich",
    ], size=PPt(17))

    s = content_slide("Infrastructure as Code & CI/CD", "Offizieller Accelerator (Deploy-Accelerator)")
    bullets(s, PInches(0.7), PInches(1.7), PInches(12), PInches(5.0), [
        "Bicep + 22 Azure Verified Modules (Microsoft-Registry)",
        "Bootstrap (Phase 0): erstellt GitHub-Repo, passwortloses OIDC & Pipelines",
        "Deploy-Accelerator liest config/inputs-github.yaml + platform-landing-zone.yaml",
        "18 geordnete Deploy-Stufen (Governance → Logging → Networking) über generierte Pipelines",
        "What-If vorab, Apply über Approval-Gate · Runbook: docs/ACCELERATOR-BOOTSTRAP.md",
    ], size=PPt(16))

    # --- Status ---
    s = content_slide("Umsetzungsstand", "Basiert auf offiziellem Microsoft-Accelerator")
    rows = [
        ("Management Groups (12) + volles ALZ-Policy-Set", "Umgesetzt"),
        ("RBAC (5 Custom-Rollen + Zuweisungen)", "Umgesetzt"),
        ("Logging via avm/ptn/alz/ama (LAW, 3 DCRs, MI)", "Umgesetzt"),
        ("Defender for Cloud (via Deploy-MDFC-* Policies)", "Umgesetzt"),
        ("Hub-Networking (FW, Bastion, GW, DDoS, DNS)", "Umgesetzt (Default an)"),
        ("Bootstrap + generierte Pipelines + OIDC", "Umgesetzt"),
        ("Identity-Ressourcen / Sentinel-Connectors", "Roadmap"),
        ("Virtual WAN (Alternative)", "Vorbereitet / optional"),
    ]
    tbl = s.shapes.add_table(len(rows)+1, 2, PInches(0.7), PInches(1.5), PInches(12), PInches(5.2)).table
    tbl.columns[0].width = PInches(8.6); tbl.columns[1].width = PInches(3.4)
    for j, htxt in enumerate(["Domäne", "Status"]):
        c = tbl.cell(0, j); c.text = htxt
        c.fill.solid(); c.fill.fore_color.rgb = PBLUE
        c.text_frame.paragraphs[0].runs[0].font.color.rgb = PWHITE
        c.text_frame.paragraphs[0].runs[0].font.bold = True
        c.text_frame.paragraphs[0].runs[0].font.size = PPt(14)
    for i, (a, b) in enumerate(rows, start=1):
        for j, val in enumerate((a, b)):
            c = tbl.cell(i, j); c.text = val
            r0 = c.text_frame.paragraphs[0].runs[0]; r0.font.size = PPt(12)
            r0.font.color.rgb = PDARK if j == 0 else (PBLUE if val == "Umgesetzt" else PGREY)
            if j == 1: r0.font.bold = True

    # --- Roadmap ---
    s = content_slide("Roadmap / nächste Schritte")
    bullets(s, PInches(0.7), PInches(1.6), PInches(12), PInches(5.2), [
        "Bootstrap im Kunden-Tenant ausführen (Deploy-Accelerator), dann What-If → Apply",
        "Netzwerk-Scope festlegen: voller Default (~€5.800/Mon.) vs. kostenarm (network_type=none / Schalter aus)",
        "Guardrail-Initiativen mit DoNotEnforce nach Compliance scharfschalten",
        "Identity-Domäne ausgestalten (Entra-Diagnose, hybride Anbindung)",
        "Sentinel-Onboarding (Data Connectors, Analyseregeln)",
        "On-Prem-Konnektivität abstimmen (VPN / ExpressRoute)",
    ], size=PPt(17))

    # --- Beratungspunkte ---
    s = content_slide("Entscheidungs- & Beratungspunkte", "Gemeinsam im Kickoff")
    bullets(s, PInches(0.7), PInches(1.7), PInches(6.0), PInches(5.0), [
        "Netzwerk-Scope & Kosten (Default ~€5.800 vs. kostenarm)",
        "Netzwerk-Topologie (Hub-Spoke vs. vWAN)",
        "IP-Adresskonzept / On-Prem-Overlap",
        "Regionen & Datenresidenz",
        "On-Prem-Anbindung: VPN vs. ExpressRoute",
        "DDoS Protection (Default an, Kostenfaktor)",
    ], size=PPt(15))
    bullets(s, PInches(6.9), PInches(1.7), PInches(5.8), PInches(5.0), [
        "RBAC-Modell & Entra-Gruppen (PIM)",
        "Guardrail-Enforcement (DoNotEnforce scharfschalten)",
        "Defender-Tier (Standard vs. Free)",
        "Sentinel-Onboarding",
        "Single- vs. Multi-Subscription",
        "Log-Aufbewahrung & Kosten",
    ], size=PPt(15))

    # --- Abschluss ---
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SW, SH, PDARK)
    add_rect(s, 0, PInches(3.0), SW, PInches(0.08), PBLUE)
    txt(s, PInches(0.8), PInches(3.15), PInches(11.7), PInches(1.0), "Vielen Dank – Fragen & Diskussion", PPt(34), PWHITE, True)
    txt(s, PInches(0.85), PInches(4.3), PInches(11.7), PInches(0.6), "Nächster Schritt: Bootstrap im Kunden-Tenant (Deploy-Accelerator), dann What-If → Apply", PPt(16), PRGB(0x9F,0xD3,0xFF))

    path = f"{OUT}/Azure-Landing-Zone-Kickoff.pptx"
    prs.save(path)
    return path

if __name__ == "__main__":
    d = build_docx()
    p = build_pptx()
    print("DOCX:", d)
    print("PPTX:", p)
