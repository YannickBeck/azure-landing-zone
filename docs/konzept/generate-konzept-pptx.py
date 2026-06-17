# -*- coding: utf-8 -*-
"""
Generiert das Bechtle-Kundenkonzept "Azure Landing Zone" als PowerPoint-Präsentation.

Basis-Vorlage : bechtle-brand/Powerpoint/Bechtle_Designvorlage_Grün_Teil1.pptx
Output        : Powerpoint/Azure-Landing-Zone-Konzept.pptx

Nutzung:
    python3 docs/konzept/generate-konzept-pptx.py

Abhängigkeit  : python-pptx  (pip install python-pptx)
"""

import os
from pptx import Presentation
from pptx.util import Pt, Inches

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
IMAGES   = os.path.join(BASE_DIR, "images")
TEMPLATE  = os.path.join(BASE_DIR, "bechtle-brand", "Powerpoint",
                         "VORLAGE_Bechtle_Praesentation.pptx")
OUTPUT    = os.path.join(BASE_DIR, "Powerpoint", "Azure-Landing-Zone-Konzept.pptx")

KUNDE = "<KUNDE>"
DATE  = "16.06.2026"

# Layout-Indices im Slide-Master
L_TITLE   = 0   # Title Slide with Picture
L_AGENDA  = 4   # Agenda
L_DIVIDER = 6   # Divider Number / Icon
L_CONTENT = 14  # Titel und Inhalt

# Beziehungs-Namespace für r:id-Attribut
_NS_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


# ─────────────────────────────────────────────────────────────────────────────
# Hilfsfunktionen
# ─────────────────────────────────────────────────────────────────────────────

def clear_slides(prs):
    """Alle bestehenden Folien aus der Vorlage entfernen."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(_NS_R + "id")
        if r_id:
            try:
                prs.part.drop_rel(r_id)
            except Exception:
                pass
        sld_id_lst.remove(sld_id)


def sl(prs, idx):
    """Slide-Layout per Index aus dem ersten Master."""
    return prs.slide_masters[0].slide_layouts[idx]


def ph(slide, idx):
    """Platzhalter per Format-Index suchen."""
    for p in slide.placeholders:
        if p.placeholder_format.idx == idx:
            return p
    return None


def set_text(slide, ph_idx, text, size=None, bold=False):
    """Einfachen Text in einen Platzhalter schreiben."""
    p = ph(slide, ph_idx)
    if p is None:
        return
    tf = p.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    if size:
        run.font.size = Pt(size)
    if bold:
        run.font.bold = True


def set_bullets(slide, ph_idx, items, size=13):
    """
    Bullet-Punkte in einen Platzhalter schreiben.
    items: Liste von str oder (str, level)-Tupeln (level 0=Haupt, 1=Unter-Bullet).
    """
    p = ph(slide, ph_idx)
    if p is None:
        return
    tf = p.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, level = (item, 0) if isinstance(item, str) else item
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = level
        run = para.add_run()
        run.text = text
        if size:
            run.font.size = Pt(size)


# ─────────────────────────────────────────────────────────────────────────────
# Folien-Bausteine
# ─────────────────────────────────────────────────────────────────────────────

def title_slide(prs):
    slide = prs.slides.add_slide(sl(prs, L_TITLE))
    set_text(slide, 0, "Azure Landing Zone")
    set_text(slide, 1, "Konzept und Umsetzungsfahrplan")
    set_text(slide, 15,
             f"{KUNDE}  ·  Bechtle GmbH & Co. KG  ·  {DATE}", size=10)
    return slide


def agenda_slide(prs):
    slide = prs.slides.add_slide(sl(prs, L_AGENDA))
    set_text(slide, 0, "Agenda")
    set_bullets(slide, 13, [
        "Management Summary",
        "1  ·  Ausgangslage und Zielbild",
        "2  ·  Methodik und Vorgehen",
        "3  ·  Zielarchitektur",
        "4  ·  Governance und Policies",
        "5  ·  Netzwerk-Architektur",
        "6  ·  Sicherheit",
        "7  ·  Monitoring und Logging",
        "8  ·  Identity und RBAC",
        "9  ·  Automatisierung und CI/CD",
        "10  ·  Kosten und Kostensteuerung",
        "11  ·  Roadmap und Phasen",
        "12  ·  Risiken und Entscheidungen",
        "13  ·  Ergänzende Randthemen",
        "14  ·  Nächste Schritte",
    ], size=13)
    return slide


def divider(prs, title, subtitle, facts=None):
    """Kapitel-Trenner-Folie (Layout 6 – Divider Number / Icon)."""
    slide = prs.slides.add_slide(sl(prs, L_DIVIDER))
    set_text(slide, 0, title)
    set_text(slide, 1, subtitle)
    if facts:
        set_bullets(slide, 13, facts, size=12)
    return slide


def content(prs, breadcrumb, title, bullets, size=13):
    """Standard-Inhalts-Folie (Layout 14 – Titel und Inhalt)."""
    slide = prs.slides.add_slide(sl(prs, L_CONTENT))
    set_text(slide, 1, breadcrumb)
    set_text(slide, 0, title)
    set_bullets(slide, 13, bullets, size=size)
    return slide


def highlight(prs, breadcrumb, statement, body=None):
    """Layout 10 – Highlight Dark: Key-Statement auf dunklem Hintergrund."""
    slide = prs.slides.add_slide(sl(prs, 10))
    set_text(slide, 1, breadcrumb, size=11)
    set_text(slide, 0, statement, size=22, bold=True)
    if body:
        set_text(slide, 13, body, size=12)
    return slide


def compare(prs, breadcrumb, title, left_title, left_items,
            right_title, right_items, size=12):
    """Layout 16 – 2x Text: Zwei-Spalten-Vergleich."""
    slide = prs.slides.add_slide(sl(prs, 16))
    set_text(slide, 1, breadcrumb)
    set_text(slide, 0, title)
    set_bullets(slide, 13, [(left_title, 0)] + [(t, 1) for t in left_items], size=size)
    set_bullets(slide, 14, [(right_title, 0)] + [(t, 1) for t in right_items], size=size)
    return slide


def diagram_slide(prs, breadcrumb, title, img_file, bullets, big=False):
    """
    Layout 19 (big=False) oder 20 (big=True) – Text links, Bild rechts.
    Nutzt den nativen PICTURE-Placeholder (PH14) statt manuell positionierter Bilder.
    """
    layout_idx = 20 if big else 19
    slide = prs.slides.add_slide(sl(prs, layout_idx))
    set_text(slide, 1, breadcrumb)
    set_text(slide, 0, title)
    if bullets:
        set_bullets(slide, 13, bullets, size=11)
    img_path = os.path.join(IMAGES, img_file)
    if os.path.exists(img_path):
        pic_ph = ph(slide, 14)
        if pic_ph is not None:
            pic_ph.insert_picture(img_path)
    return slide


def picture_slide(prs, title, filename, callouts):
    """
    Folie: Original-Diagramm links (66 %), Kommentarboxen rechts (30 %).
    callouts: Liste von (header, text)-Tupeln – max. 4 sinnvoll.
    Das Diagramm bleibt vollständig unverändert; alle Bechtle-Infos
    stehen als farbige Boxen in der rechten Spalte.
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    from pptx.oxml.ns import qn as _qn
    from lxml import etree

    G_DARK  = RGBColor(0x07, 0x50, 0x33)   # Bechtle Dunkelgrün
    G_MID   = RGBColor(0x23, 0xA9, 0x6A)   # Bechtle Grün
    G_BG    = RGBColor(0xF2, 0xF9, 0xF5)   # sehr helles Grün (Box-Hintergrund)
    GREY    = RGBColor(0x59, 0x59, 0x59)

    slide = prs.slides.add_slide(sl(prs, 39))  # Leer/Blank

    # ── Titelleiste oben ──────────────────────────────────────────────────────
    tb = slide.shapes.add_textbox(Inches(0.15), Inches(0.08), Inches(13.0), Inches(0.48))
    tf = tb.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = G_DARK

    # Trennlinie unter dem Titel (dünnes grünes Rechteck)
    line = slide.shapes.add_shape(1, Inches(0.15), Inches(0.60), Inches(13.0), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = G_MID
    line.line.fill.background()

    # ── Originaldiagramm links ────────────────────────────────────────────────
    img_path = os.path.join(IMAGES, filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path,
            left=Inches(0.15),
            top=Inches(0.68),
            width=Inches(8.75)   # lässt rechts 4,35" für Kommentarboxen
        )

    # Quellenangabe links unten
    src_tb = slide.shapes.add_textbox(Inches(0.15), Inches(7.25), Inches(8.75), Inches(0.22))
    src_tf = src_tb.text_frame
    src_r = src_tf.paragraphs[0].add_run()
    src_r.text = "Quelle: Microsoft Cloud Adoption Framework · learn.microsoft.com (unverändert)"
    src_r.font.size = Pt(7)
    src_r.font.color.rgb = GREY

    # ── Kommentarboxen rechts ─────────────────────────────────────────────────
    BOX_LEFT  = Inches(9.15)
    BOX_W     = Inches(4.00)
    AREA_TOP  = Inches(0.68)
    AREA_H    = Inches(6.60)
    n = len(callouts)
    box_h = AREA_H / n - Inches(0.08)   # kleiner Abstand zwischen Boxen

    for i, (header, body) in enumerate(callouts):
        bx_top = AREA_TOP + i * (AREA_H / n)

        # Box-Hintergrund
        box = slide.shapes.add_shape(1, BOX_LEFT, bx_top, BOX_W, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = G_BG
        box.line.color.rgb = G_DARK
        box.line.width = Pt(1.25)

        # Grüner Akzentstreifen links an der Box
        accent = slide.shapes.add_shape(
            1, BOX_LEFT, bx_top, Inches(0.07), box_h)
        accent.fill.solid()
        accent.fill.fore_color.rgb = G_DARK
        accent.line.fill.background()

        # Text in der Box
        txt_box = slide.shapes.add_textbox(
            BOX_LEFT + Inches(0.12), bx_top + Inches(0.07),
            BOX_W - Inches(0.18), box_h - Inches(0.10))
        tf2 = txt_box.text_frame
        tf2.word_wrap = True

        p_hdr = tf2.paragraphs[0]
        r_hdr = p_hdr.add_run()
        r_hdr.text = header
        r_hdr.font.bold = True
        r_hdr.font.size = Pt(9)
        r_hdr.font.color.rgb = G_DARK

        for line_txt in body.split("\n"):
            p_b = tf2.add_paragraph()
            r_b = p_b.add_run()
            r_b.text = line_txt
            r_b.font.size = Pt(8)
            r_b.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Hauptfunktion
# ─────────────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    # ── 1. Titelfolie ────────────────────────────────────────────────────────
    title_slide(prs)

    # ── 2. Agenda ────────────────────────────────────────────────────────────
    agenda_slide(prs)

    # ── Management Summary ───────────────────────────────────────────────────
    highlight(prs,
        "Management Summary",
        "Governance · Sicherheit · Netzwerk\nals Code – reproduzierbar und auditfähig.",
        "Bechtle setzt auf den Microsoft ALZ Bicep Accelerator: "
        "149 Policies · 12 MGs · 18 Deploy-Stufen · OIDC · ~€1.050/Monat"
    )
    divider(prs,
        "Management Summary",
        "Governance · Netzwerk · Sicherheit · Automatisierung – als IaC-Fundament",
        [
            "12 Management Groups",
            "149 Policy-Definitionen",
            "42 Initiativen  ·  123 Assignments",
            "22 AVM-Module",
            "18 Deploy-Stufen",
            "~€5.800/Monat (vollständig steuerbar)",
        ]
    )
    content(prs, "Management Summary", "Die vier Kernnutzen der Azure Landing Zone", [
        ("Governance & Policies: 149 Policy-Defs, 42 Initiativen, 123 Assignments – "
         "automatische Guardrails auf allen MG-Ebenen", 0),
        ("Netzwerk-Hub (Hub-and-Spoke): Azure Firewall, Bastion, Private DNS – "
         "sichere Konnektivität für alle Workloads", 0),
        ("Sicherheit & Defender: Microsoft Defender for Cloud automatisch via Policy aktiviert – "
         "keine manuelle Konfiguration", 0),
        ("IaC & Automatisierung: 18 Deploy-Stufen, OIDC (passwortlos), "
         "What-If vor Apply, Approval-Gate", 0),
        ("Kosten: Microsoft-Default ~€5.800/Monat  ↔  Pilot mit network_type: none ≈ €0", 0),
        ("Empfehlung: Pilot mit 1 Subscription, network_type: none, "
         "Policies im DoNotEnforce-Modus – Governance vom ersten Tag an", 0),
    ])

    # ── 1. Ausgangslage & Zielbild ───────────────────────────────────────────
    divider(prs,
        "1  ·  Ausgangslage und Zielbild",
        "Von der ersten Subscription zur skalierbaren Cloud-Grundstruktur"
    )
    content(prs, "Ausgangslage und Zielbild", "Warum eine Azure Landing Zone?", [
        ("Ausgangssituation: Azure-Einstieg ohne Governance-Modell – "
         "Ressourcen entstehen ad hoc, ohne einheitliche Leitplanken", 0),
        ("Chance: Von Anfang an auf industrieerprobter Basis aufbauen – "
         "statt technische Schulden durch spätere Nachrüstung", 0),
        ("Zielbild: Standardisierte, mandantenweite Grundstruktur als stabiles Fundament "
         "für alle künftigen Workloads", 0),
        ("12 MGs als strukturierter Verwaltungsrahmen für alle Subscriptions", 1),
        ("Vollständiges ALZ-Policy-Set für automatische Governance ohne manuelle Eingriffe", 1),
        ("Reproduzierbares IaC-Modell: audit-fähig, rückrollbar, nachvollziehbar", 1),
        ("Abgrenzung: ALZ = Fundament – Migration, Anwendungsentwicklung und "
         "Applikationsbetrieb sind eigenständige Folgeprojekte", 0),
    ])

    # ── 2. Methodik & Vorgehen ────────────────────────────────────────────────
    divider(prs,
        "2  ·  Methodik und Vorgehen",
        "Microsoft ALZ Bicep Accelerator – kein Eigenbau, sondern Microsoft-Standard"
    )
    content(prs, "Methodik und Vorgehen",
            "ALZ Bicep Accelerator + 18 Deployment-Stufen", [
        ("Microsoft-Standard: derselbe Accelerator, den Microsoft für tausende "
         "Enterprise-Kunden empfiehlt und pflegt", 0),
        ("Policy-Set komplett: 149 Definitionen, 42 Initiativen, 123 Assignments – "
         "automatisch aus MCR (Microsoft Container Registry)", 0),
        ("22 Azure Verified Modules (AVM): Microsoft-zertifizierte Bicep-Module, "
         "kein Eigenbau-IaC", 0),
        ("18 geordnete Deploy-Stufen: Governance (1–12) → RBAC (13–15) → "
         "Logging (16) → Netzwerk (17/18)", 0),
        ("OIDC Federated Identity: passwortlose Pipelines – keine Secrets oder "
         "Zertifikate", 0),
        ("What-If vor Apply: Pipeline zeigt geplante Änderungen zur Freigabe, "
         "bevor etwas ausgeführt wird", 0),
        ("Qualitätssicherung: bicep build aller 20 Templates offline; "
         "20/20 grün (Bicep 0.44.1)", 0),
    ])

    # ── 3. Zielarchitektur ────────────────────────────────────────────────────
    divider(prs,
        "3  ·  Zielarchitektur",
        "12 Management Groups + Hub-and-Spoke-Netzwerk in zwei Regionen"
    )
    content(prs, "Zielarchitektur", "Management Group Hierarchie (12 MGs)", [
        ("Tenant Root", 0),
        ("alz (Intermediate Root) – alle ALZ-Policies erben von hier", 1),
        ("alz-platform: connectivity · identity · management · security", 1),
        ("alz-platform-connectivity  → Hub-Netzwerk, Firewall, DNS", 1),
        ("alz-platform-identity      → Identity-Dienste (Entra, AD DS)", 1),
        ("alz-platform-management    → Logging, Monitoring, Operations", 1),
        ("alz-platform-security      → Security-Tooling, Defender", 1),
        ("alz-landingzones: corp · online · local", 1),
        ("alz-sandbox  ·  alz-decommissioned", 1),
        ("Region: Germany West Central (primär 10.0.0.0/22) + "
         "North Europe (sekundär 10.1.0.0/22)", 0),
        ("Subscription-Strategie: Single-Sub für Pilot → "
         "4 dedizierte Platform-Subs für Produktion", 0),
    ])
    diagram_slide(prs,
        "Zielarchitektur",
        "Hub-and-Spoke-Topologie",
        "alz-hub-spoke.png",
        [
            "Bechtle-Empfehlung: GWC · Firewall Standard · ~€1.050/Mon.",
            "Connectivity Sub: Firewall · Bastion · VPN GW · DNS Resolver",
            "Landing Zones: alz-corp · alz-online · alz-sandbox",
            "On-Prem: VPN aktiv · ExpressRoute deferred",
            "SIEM: CrowdStrike via Event-Hub-Export",
        ],
        big=True
    )
    diagram_slide(prs,
        "Zielarchitektur",
        "Management Group Hierarchie (12 MGs)",
        "alz-mg-hierarchy.png",
        [
            "Vollständig per ALZ Bicep Accelerator deployt",
            "149 Definitionen · 42 Initiativen · 123 Assignments",
            "DoNotEnforce: Audit-Modus schützt Bestandsressourcen",
            "What-If vor jedem Subscription-Move",
            "Exemptions für Legacy-Ressourcen möglich",
        ],
        big=False
    )

    # ── 4. Governance & Policies ──────────────────────────────────────────────
    divider(prs,
        "4  ·  Governance und Policies",
        "149 Policy-Definitionen · 42 Initiativen · 123 Assignments",
        [
            "9 MG-Ebenen mit Assignments",
            "5 Custom RBAC-Rollen",
            "DoNotEnforce → sanftes Onboarding",
        ]
    )
    content(prs, "Governance und Policies", "Das vollständige ALZ-Policy-Set", [
        ("Policy-Definitionen (Custom): 149 – Monitoring (55), Netzwerk (20), "
         "Storage (16), SQL (13), Guardrails je Dienst", 0),
        ("Policy-Set-Definitionen (Initiativen): 42 – Deploy-MDFC-Config, "
         "Deploy-Private-DNS-Zones, Enforce-Guardrails-* je Dienst", 0),
        ("Policy-Assignments: 123 – verteilt auf 9 MG-Ebenen; "
         "Kind-MGs erben zusätzlich", 0),
        ("alz-Root: 17  ·  alz-landingzones: 53  ·  alz-platform: 40  ·  "
         "alz-corp: 5  ·  weitere: 8", 1),
        ("5 Custom RBAC-Rollen: Subscription-Owner · Security-Ops · "
         "Network-Mgmt · App-Owners · Subnet-Contributor", 0),
        ("DoNotEnforce-Modus: Policies greifen (Compliance-Dashboard), "
         "blockieren aber nichts → sanftes Onboarding", 0),
        ("Stufenplan: Audit → selektiv Enabled → "
         "vollständiges Enforcement nach Remediation", 0),
    ])

    # ── 5. Netzwerk-Architektur ───────────────────────────────────────────────
    divider(prs,
        "5  ·  Netzwerk-Architektur",
        "Hub-and-Spoke · Bechtle-Empfehlung: 1 Region · Firewall Standard · ~€1.050/Monat"
    )
    content(prs, "Netzwerk-Architektur", "Hub-and-Spoke: Bechtle-Empfehlung vs. Microsoft-Default", [
        ("Bechtle-Empfehlung: 1× Hub-VNet GWC 10.0.0.0/22 – alle Kernfunktionen, "
         "eine Region, ~€1.050/Monat", 0),
        ("Microsoft-Default: 2× Hub-VNets (GWC + NE) – alle Dienste, beide Regionen, "
         "~€5.800/Monat", 0),
        ("Azure Firewall Standard (1×): ~€700/Monat  [Default: Premium 2× = ~€2.200]", 0),
        ("Azure Bastion Standard (1×): ~€120/Monat  [Default: 2× = €240]", 0),
        ("VPN Gateway VpnGw1AZ (1×): ~€140/Monat  [Default: 2× = €280]", 0),
        ("ExpressRoute Gateway: deaktiviert – bei ER-Leitung aktivieren [Default: 2× = €560]", 0),
        ("DDoS Protection: deaktiviert – bei internet-facing Workloads aktivieren [Default: €2.500]", 0),
        ("Private DNS Zones + DNS Resolver: aktiv – unveraendert  ~€40/Monat", 0),
    ])
    diagram_slide(prs,
        "Netzwerk-Architektur",
        "Connectivity Subscription im Detail",
        "alz-hub-spoke.png",
        [
            "Azure Firewall Standard: FQDN · NAT · Threat Intel · ~€700/Mon.",
            "VPN GW VpnGw1AZ: ~€140/Mon.  |  Bastion: RDP/SSH ohne Public IP · ~€120",
            "Key Vault: Zertifikate für Monitoring-TLS",
            "Event Hub → CrowdStrike SIEM (Log-Export)",
            "UDR: aller Spoke-Traffic zwingend über Firewall",
        ],
        big=False
    )
    compare(prs,
        "Netzwerk-Architektur",
        "Microsoft-Default vs. Bechtle-Empfehlung",
        "Microsoft-Default  ~€5.800/Mon.",
        [
            "Azure Firewall Premium 2× · ~€2.200",
            "DDoS Network Protection · ~€2.500",
            "ExpressRoute GW 2× · ~€560",
            "VPN GW 2× · ~€280",
            "Azure Bastion 2× · ~€240",
            "2 Regionen: GWC + North Europe",
        ],
        "Bechtle-Empfehlung  ~€1.050/Mon.",
        [
            "Azure Firewall Standard 1× · ~€700",
            "DDoS: deferred (bei Bedarf aktivieren)",
            "ExpressRoute GW: deferred",
            "VPN GW 1× · ~€140",
            "Azure Bastion 1× · ~€120",
            "1 Region: Germany West Central",
        ]
    )

    # ── 6. Sicherheit ─────────────────────────────────────────────────────────
    divider(prs,
        "6  ·  Sicherheit",
        "Defender for Cloud · Netzwerksegmentierung · Automatische Remediation"
    )
    content(prs, "Sicherheit", "Mehrschichtige Sicherheitsarchitektur", [
        ("Präventiv (Policies): Deny-Subnet-Without-Nsg, Deny-MgmtPorts-Internet, "
         "Deny-Public-IP, Deny-Public-Endpoints (Corp)", 0),
        ("Defender for Cloud: auto-aktiviert via Deploy-MDFC-Config-H224 auf "
         "alz-Root (DeployIfNotExists)", 0),
        ("Defender for Endpoint: Deploy-MDEndpoints auf allen VMs – "
         "im Defender Server-Plan enthalten", 0),
        ("MDFC Benchmarks: Microsoft Cloud Security Benchmark v1 + v2 "
         "(kostenlos, Audit)", 0),
        ("Activity Logs + Diagnose: alle Subscriptions und Ressourcen "
         "→ Log Analytics (via DINE-Policies)", 0),
        ("Backup: Deploy-VM-Backup sichert VMs ohne Backup-Tag automatisch "
         "in Recovery Vault", 0),
        ("Roadmap: Microsoft Sentinel als SIEM-Folgestufe "
         "(nicht Bestandteil der Landing Zone)", 0),
    ])

    # ── 7. Monitoring & Logging ───────────────────────────────────────────────
    divider(prs,
        "7  ·  Monitoring und Logging",
        "Zentraler Log Analytics Workspace + 3 Data Collection Rules"
    )
    content(prs, "Monitoring und Logging",
            "Zentrales Logging (AVM-Pattern avm/ptn/alz/ama:0.2.0)", [
        ("Log Analytics Workspace: law-alz-<region>, PerGB2018, "
         "365 Tage Retention", 0),
        ("3 Data Collection Rules: VM Insights · Change Tracking · "
         "Defender for SQL", 0),
        ("User-Assigned Managed Identity: mi-alz-<region> für passwortlose "
         "Policy-Remediation (DeployIfNotExists)", 0),
        ("DINE-Policies: jede neue VM verbindet sich automatisch mit LAW + DCRs – "
         "keine manuelle Konfiguration", 0),
        ("Deploy-Diag-LogsCat: leitet alle Diagnose-Kategorien aller Ressourcen "
         "in Log Analytics", 0),
        ("Kosten: 5 GB/Monat kostenlos; danach ~€2/GB; "
         "Automation Account optional (default: aus)", 0),
    ])

    # ── 8. Identity & RBAC ────────────────────────────────────────────────────
    divider(prs,
        "8  ·  Identity und RBAC",
        "5 Custom-Rollen + Entra-ID-Gruppen-basierte Zuweisungen"
    )
    content(prs, "Identity und RBAC", "5 Custom RBAC-Rollen", [
        ("Subscription-Owner (alz): delegierter Owner, eingeschränkt ohne volle "
         "Owner-Rechte (1 Action)  → Plattform-Admins", 0),
        ("Security-Operations (alz): horizontale Sicherheitssicht über alle "
         "Subscriptions (12 Actions)  → Security-Team", 0),
        ("Network-Management (alz): VNets, UDRs, NSGs, NVAs plattformweit "
         "(4 Actions)  → Netzwerk-Team", 0),
        ("Application-Owners (alz): Contributor auf Resource-Group-Ebene "
         "(1 Action)  → App-/Ops-Teams", 0),
        ("Network-Subnet-Contributor (alz): vollständiger Subnetz-Zugriff "
         "(8 Actions)  → Netzwerk-Ops", 0),
        ("Zuweisung: Entra-ID-Gruppen-Object-IDs in Bicep-Params; "
         "leeres Array = No-Op (gefahrlos)", 0),
        ("Roadmap: PIM, Conditional Access, hybride Identitäten als "
         "Folge-Baustein (Identity-Baseline)", 0),
    ])

    # ── 9. Automatisierung & CI/CD ────────────────────────────────────────────
    divider(prs,
        "9  ·  Automatisierung und CI/CD",
        "Bootstrap → 18 Deploy-Stufen → GitOps-Betrieb"
    )
    content(prs, "Automatisierung und CI/CD",
            "Pipeline-Workflow (GitHub Actions)", [
        ("Bootstrap (Phase 0, Terraform): Resource Group + Storage (State) + "
         "Managed Identity (OIDC Federated Credentials)", 0),
        ("GitHub-Repo + Environments + Approval-Gates automatisch generiert", 0),
        ("5 Pipeline-Schritte: Trigger → Validation (bicep build) → "
         "What-If → Approval → Apply", 0),
        ("18 Stufen sequenziell: Governance-MGs (1–12) → RBAC (13–15) → "
         "Logging (16) → Hub-Netzwerk (17) oder vWAN (18)", 0),
        ("OIDC Federated Identity: keine Passwörter, keine Zertifikate, "
         "keine Secrets in Pipelines", 0),
        ("GitOps: PR → Review → Approval → Deploy; "
         "jede Änderung mit Autor + Timestamp in Git-History", 0),
        ("Rückrollbarkeit: neuer Commit → Rollback-Pipeline; "
         "vollständiges Audit-Trail in Azure Activity Log", 0),
    ])

    # ── 10. Kosten & Sizing ───────────────────────────────────────────────────
    highlight(prs,
        "Kosten und Kostensteuerung",
        "~€1.050 / Monat",
        "Bechtle-Empfehlung: voller Funktionsumfang, eine Region, ALZ-konform – "
        "4× günstiger als Microsoft-Default (~€5.800)"
    )
    divider(prs,
        "10  ·  Kosten und Kostensteuerung",
        "Bechtle-Empfehlung: ~€1.050/Monat – voller Funktionsumfang, eine Region"
    )
    content(prs, "Kosten und Kostensteuerung",
            "Gestufter Ausbau: €0 → €1.050 → €5.800", [
        ("Stufe 1 – Governance-Pilot: network_type: none → "
         "MGs + volles Policy-Set + Logging  ca. €50/Monat (nur LAW)", 0),
        ("Stufe 2 – Bechtle-Empfehlung: Firewall Standard, 1 Region, "
         "kein DDoS/ER  →  ~€1.050/Monat", 0),
        ("  Firewall Standard 1×: ~€700  |  Bastion 1×: ~€120  |  VPN GW 1×: ~€140  "
         "|  DNS: ~€40  |  LAW: ~€50", 1),
        ("Stufe 3 – Geo-Redundanz: + Sekundar-Hub NE  →  ~€1.800/Monat", 0),
        ("Stufe 4 – Firewall Premium: IDPS / TLS / URL-Filterung  →  ~€2.400/Monat", 0),
        ("Stufe 5 – DDoS Protection: bei internet-facing Workloads  →  ~€4.900/Monat", 0),
        ("Stufe 6 – Microsoft-Default: alle Dienste, beide Regionen  →  ~€5.800/Monat", 0),
        ("Alle Stufen sind additive Aenderungen – kein Rueckbau erforderlich", 0),
    ])

    content(prs, "Kosten und Kostensteuerung",
            "Konfigurationsvergleich: Vorteile und Nachteile (1/2)", [
        ("Pilot  ~€50/Monat  (network_type: none)", 0),
        ("+ Nullrisiko; volles Governance-Set (MGs, Policies, Logging)", 1),
        ("– Keine Netzwerkkonnektivität; Workloads nicht deploybar", 1),
        ("  Empfehlung: Einstieg, PoC, Compliance-Audit ohne Netzwerkbedarf", 1),
        ("Bechtle-Empfehlung  ~€1.050/Monat  (Standard FW, 1 Region)", 0),
        ("+ Voller Funktionsumfang; 4× günstiger als Microsoft-Default", 1),
        ("+ Firewall, Bastion, VPN, DNS aktiv; In-Place-Upgrade auf Premium möglich", 1),
        ("– Keine zweite Region; kein IDPS/TLS; kein DDoS-Schutz", 1),
        ("  Empfehlung: Standard-Produktionsumgebungen, erste Workloads", 1),
        ("Geo-Redundanz  ~€1.800/Monat  (+ Hub North Europe)", 0),
        ("+ Regionale Ausfallsicherheit; zwei unabhängige Netzwerkhubs", 1),
        ("– Doppelte Netzwerkkosten; erhöhte Betriebskomplexität", 1),
        ("  Empfehlung: SLA-kritische Systeme, Multi-Region-Anforderungen", 1),
    ], size=11)

    content(prs, "Kosten und Kostensteuerung",
            "Konfigurationsvergleich: Vorteile und Nachteile (2/2)", [
        ("Firewall Premium  ~€2.400/Monat  (+ IDPS / TLS-Inspektion)", 0),
        ("+ IDPS, TLS-Inspektion, URL-Filterung; In-Place-Upgrade ohne Rebuild", 1),
        ("– +€1.350 ggü. Bechtle-Empfehlung; Zertifikat-Rollout erforderlich", 1),
        ("  Empfehlung: Hohe Sicherheitsanforderungen, regulierte Branchen", 1),
        ("DDoS Protection  ~€4.900/Monat  (+ Network Protection Plan)", 0),
        ("+ Vollständiger DDoS-Schutz für alle öffentlichen IPs; SLA-Garantie", 1),
        ("– ~€2.500/Monat Fixkosten (Plan Fee), unabhängig von Angriffen", 1),
        ("  Empfehlung: Nur bei internet-facing Workloads mit realem DDoS-Risiko", 1),
        ("Microsoft-Default  ~€5.800/Monat  (alle Dienste, 2 Regionen)", 0),
        ("+ Maximale Redundanz; Enterprise-Standard; Herstellerempfehlung", 1),
        ("– Höchste Kosten; viele Komponenten initial ungenutzt; Overkill für Start", 1),
        ("  Empfehlung: Unternehmenskritische Systeme mit strengster Compliance", 1),
    ], size=11)

    content(prs, "Kosten und Kostensteuerung",
            "Bechtle-Varianten: Vier Konfigurationen im Vergleich", [
        ("Option A – Bechtle-Standard  ~€1.050/Monat  (Baseline)", 0),
        ("FW Standard 1× ~€700 | VPN GW 1× ~€140 | Bastion ~€120 | DNS ~€25 | LAW ~€50", 1),
        ("ALZ-konform: Ja  |  Differenz: Referenz", 1),
        ("Option B – Bechtle-Optimiert  ~€770/Monat  (ALZ-konform)", 0),
        ("FW Standard reserved ~€560 | VPN Gateway deferred (€0) | Bastion ~€120 | DNS ~€25 | LAW ~€50", 1),
        ("ALZ-konform: Ja  |  Differenz: −€280/Mon. (26 % günstiger)", 1),
        ("Option C – Bechtle-Budget  ~€500/Monat  (⚠ ALZ-Bruch)", 0),
        ("FW Basic ~€300 | VPN Gateway deferred (€0) | Bastion ~€120 | DNS ~€25 | LAW ~€50", 1),
        ("ALZ-konform: NEIN  |  Differenz: −€550/Mon. (52 % günstiger)", 1),
        ("Bruch: Firewall Basic = keine App-Rules, kein Threat Intel, kein Autoscaling", 1),
        ("Option D – Microsoft-Default  ~€5.800/Monat  (maximaler Standard)", 0),
        ("FW Premium 2× ~€2.200 | DDoS Plan ~€2.500 | ER GW 2× ~€560 | VPN 2× | Bastion 2× | 2 Regionen", 1),
        ("ALZ-konform: Ja (vollständig)  |  Differenz: +€4.750/Mon. (452 % teurer als A)", 1),
    ], size=11)

    # ── 11. Roadmap & Phasen ──────────────────────────────────────────────────
    divider(prs,
        "11  ·  Roadmap und Phasen",
        "6 Phasen – von Kickoff bis laufendem Betrieb"
    )
    content(prs, "Roadmap und Phasen", "6 Projektphasen im Überblick", [
        ("Phase 1 – Beratung & Kickoff: Discovery-Workshop; "
         "Entscheidungsprotokoll (Region, Netzwerk, DDoS, Enforcement)", 0),
        ("Phase 2 – Bootstrap & Grundgerüst: Deploy-Accelerator; "
         "MGs + Policies (DoNotEnforce); Logging aktiv  ≈ €0", 0),
        ("Phase 3 – Netzwerk-Pilot: Hub-VNets + DNS; erste Spoke-VNets; "
         "ggf. minimale Firewall  €15–€1.300/Monat", 0),
        ("Phase 4 – Sicherheit & Enforcement: Defender aktivieren; "
         "Guardrails einschalten; Remediation  + Defender-Kosten", 0),
        ("Phase 5 – Workload-Onboarding: erste Applikation in ALZ; "
         "Spoke-Netzwerke; RBAC-Gruppen befüllen", 0),
        ("Phase 6 – Betrieb & Optimierung: Cost Management; Tagging; "
         "Monitoring-Dashboards; Policy-Weiterentwicklung", 0),
    ])

    # ── 12. Risiken & Entscheidungen ──────────────────────────────────────────
    divider(prs,
        "12  ·  Risiken und Entscheidungen",
        "Klare Entscheidungen für einen sicheren Start"
    )
    content(prs, "Risiken und Entscheidungen",
            "Top-Risiken und Entscheidungsbedarf", [
        ("DDoS im Default AN → ~€2.500/Monat unerwartet; "
         "deployDdosProtectionPlan: false für Start", 0),
        ("Subscription-Typ: kein Free Tier → "
         "EA/MCA/PAYG vor Bootstrap sicherstellen", 0),
        ("IP-Adresskonflikt On-Prem mit Hub-VNets → "
         "IP-Plan vor Netzwerk-Deploy abstimmen", 0),
        ("Defender-Pläne auto. per Policy aktiviert → "
         "Tier-Konfiguration vor Apply reviewen", 0),
        ("Entscheidung: network_type: none (€0) vs. hubNetworking (~€5.800) "
         "für Pilot", 0),
        ("Entscheidung: Enforcement-Zeitpunkt – wann von DoNotEnforce auf "
         "Enabled (Empfehlung: 4–8 Wochen Audit-Phase)", 0),
        ("Entscheidung: On-Prem-Anbindung – "
         "VPN vs. ExpressRoute vs. keins (je nach SLA und Bandbreite)", 0),
    ])

    # ── 13. Ergänzende Randthemen ─────────────────────────────────────────────
    divider(prs,
        "13  ·  Ergänzende Randthemen",
        "Außerhalb des ALZ-Scopes – aber eng verknüpft"
    )
    content(prs, "Ergänzende Randthemen",
            "Abhängigkeiten und Folge-Bausteine", [
        ("Entra ID & Identity: PIM, Conditional Access, hybride Identitäten → "
         "Identity-Baseline als Folgestufe", 0),
        ("On-Prem-Anbindung: VPN-Gateway (deployt, aktiv) oder ExpressRoute → "
         "Leitungsbestellung + BGP-Konfiguration außerhalb Scope", 0),
        ("Backup & DR: Deploy-VM-Backup aktiv via Policy; "
         "Recovery Vault-Konfiguration und DR-Tests als Folgestufe", 0),
        ("Kostenmanagement: Azure Cost Management Dashboards, "
         "Budget-Alerts je MG/Sub, FinOps-Prozesse", 0),
        ("Compliance & Regulatorik: BSI, ISO 27001, DSGVO → "
         "Defender for Cloud Regulatory Compliance Dashboard", 0),
        ("Workload-Migration: Azure Migrate Assessment → "
         "ALZ ist das Ziel, nicht der Migrations-Prozess selbst", 0),
    ])

    # ── 14. Nächste Schritte ──────────────────────────────────────────────────
    divider(prs,
        "14  ·  Nächste Schritte",
        "Empfohlener Pilot-Pfad in 6 Schritten – ohne Kostenrisiko"
    )
    content(prs, "Nächste Schritte",
            "Sofortmaßnahmen und empfohlener Pilot-Pfad", [
        ("1.  Subscription-ID + GitHub PAT in "
         "config/inputs-github.yaml eintragen", 0),
        ("2.  Bootstrap: Deploy-Accelerator -inputs config/inputs-github.yaml "
         "config/platform-landing-zone.yaml", 0),
        ("3.  Pipeline mit What-If prüfen → "
         "welche Ressourcen werden erstellt?", 0),
        ("4.  Apply genehmigen – network_type: none (≈€0) für "
         "risikofreien Pilot-Start", 0),
        ("5.  Policy-Compliance-Dashboard reviewen → "
         "welche Guardrails greifen, welche nicht?", 0),
        ("6.  Schrittweise Enforcement und Netzwerk-Dienste aktivieren "
         "nach Audit-Phase", 0),
        ("Bechtle begleitet: Kickoff-Workshop → Bootstrap-Begleitung → "
         "Pilot → Workload-Onboarding", 0),
    ])

    # ── Abschluss-Folie ───────────────────────────────────────────────────────
    divider(prs,
        "Vielen Dank",
        "Nächster Schritt: Kickoff-Workshop vereinbaren",
        [
            KUNDE,
            " ",
            "Bechtle GmbH & Co. KG",
            "Gottlieb-Daimler-Straße 2 · 68165 Mannheim",
            "T +49 621 87503-0 · www.bechtle.com",
        ]
    )

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    slide_count = len(prs.slides)
    print(f"Gespeichert: {OUTPUT}")
    print(f"Folien: {slide_count}")


if __name__ == "__main__":
    build()
